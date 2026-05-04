"""
Vigilent Client Walkthrough Deck Generator
===========================================
Auto-generates a 2-slide PPTX from the scoring model so numbers always
match the code.  Follows Ayush's format: step-by-step worked examples
with real data center values.

Usage:
    python3 generate_client_deck.py

Output:
    output/vigilent_client_walkthrough.pptx
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from vigilent_engine import compute_score, SCORING_CONFIG

# ═══════════════════════════════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════════════════════════════

OUTPUT_PATH = "output/vigilent_client_walkthrough.pptx"
SCORED_JSON = "output/scored_datacenters.json"
INPUTS_JSON = "inputs_spec.json"

# Example DC for worked calculation (SDC Manhattan, 18 MW, Excellent)
EXAMPLE_DC_NAME = "SDC Manhattan"

# Branding colors
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
MEDIUM_BLUE = RGBColor(0x30, 0x3F, 0x9F)
LIGHT_BLUE = RGBColor(0xE8, 0xEA, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x42, 0x42, 0x42)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)

# Slide dimensions (widescreen 13.333 x 7.5 in)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=12, bold=False, color=DARK_GRAY,
             space_before=Pt(4), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = space_before
    return p


def add_table(slide, left, top, width, height, headers, rows,
              header_color=DARK_BLUE, font_size=10):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternate row shading
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: MODEL OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_1(prs, all_dcs, inputs_spec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Title bar
    add_textbox(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.8),
                "  Vigilent Optimization Model: Client Walkthrough",
                font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    # Title bar background
    title_shape = slide.shapes[0]
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE

    # --- Left column: Model description + Assumptions ---
    left_x = Inches(0.4)
    col_w = Inches(6.0)

    tf = add_textbox(slide, left_x, Inches(1.0), col_w, Inches(0.5),
                     "How the Model Works", font_size=16, bold=True, color=DARK_BLUE)

    tf = add_textbox(slide, left_x, Inches(1.5), col_w, Inches(1.8),
                     "The Vigilent composite score ranks data centers on a 0\u2013100 scale "
                     "across five factors, weighted by business impact. Each factor is "
                     "normalized to 0\u2013100, then combined via weighted sum. Higher score = "
                     "better target for Vigilent deployment.",
                     font_size=11, color=DARK_GRAY)

    # Classification tiers
    tier_counts = {}
    for dc in all_dcs:
        c = dc["classification"]
        tier_counts[c] = tier_counts.get(c, 0) + 1

    tf = add_textbox(slide, left_x, Inches(2.8), col_w, Inches(0.4),
                     "Classification Tiers", font_size=14, bold=True, color=DARK_BLUE)

    tiers = [
        ("Excellent (75\u2013100)", tier_counts.get("Excellent", 0)),
        ("Good (50\u201374)", tier_counts.get("Good", 0)),
        ("Moderate (25\u201349)", tier_counts.get("Moderate", 0)),
        ("Low (0\u201324)", tier_counts.get("Low", 0)),
    ]
    tier_text = "    ".join(f"{name}: {count} DCs" for name, count in tiers)
    add_textbox(slide, left_x, Inches(3.2), col_w, Inches(0.4),
                tier_text, font_size=11, color=DARK_GRAY)

    # Assumptions
    tf = add_textbox(slide, left_x, Inches(3.8), col_w, Inches(0.4),
                     "Key Assumptions & Sources", font_size=14, bold=True, color=DARK_BLUE)

    assumptions = [
        "Baseline PUE: 1.55 (Uptime Institute 2023 global average)",
        "Load Growth Rate: 10% (industry consensus estimate)",
        "Energy % of OPEX: 40% (Gartner / McKinsey benchmark)",
        "Vigilent Investment: $1,500,000 per deployment",
        "Energy Reduction: 10% with Vigilent (standard offering)",
        "Water Reduction: 5% with Vigilent (standard offering)",
        "Electricity Rates: EIA state-level commercial rates ($/kWh)",
    ]
    tf = add_textbox(slide, left_x, Inches(4.2), col_w, Inches(2.8),
                     assumptions[0], font_size=10, color=DARK_GRAY)
    for a in assumptions[1:]:
        add_para(tf, a, font_size=10, color=DARK_GRAY, space_before=Pt(2))

    add_para(tf, "", font_size=6)
    add_para(tf, "3 of 9 inputs are industry estimates (PUE, load growth, OPEX %).",
             font_size=10, bold=True, color=ORANGE, space_before=Pt(6))
    add_para(tf, "Site-specific data from the client will improve accuracy.",
             font_size=10, color=ORANGE, space_before=Pt(1))

    # --- Right column: Factor scoring table ---
    right_x = Inches(6.8)
    right_w = Inches(6.0)

    add_textbox(slide, right_x, Inches(1.0), right_w, Inches(0.5),
                "Scoring Factors", font_size=16, bold=True, color=DARK_BLUE)

    factor_headers = ["Factor", "Weight", "Direction", "Min", "Max", "Formula"]
    factor_display = {
        "savings_per_mw": ("Savings per MW", "35%", "Higher = Better", "$0", "$300,000",
                           "annual_savings / dc_size_mw"),
        "payback_period": ("Payback Period", "25%", "INVERTED\n(Lower = Better)", "0 yr", "5 yr",
                           "investment_cost / annual_savings"),
        "impact_on_opex": ("OPEX Impact", "20%", "Higher = Better", "0%", "10%",
                           "energy_reduction % \u00d7 energy_pct_opex"),
        "water_savings_pct": ("Water Savings", "10%", "Higher = Better", "0%", "8%",
                              "water_reduction_pct (Vigilent param)"),
        "load_growth": ("Load Growth", "10%", "Higher = Better", "0%", "15%",
                        "load_growth_rate (DC param)"),
    }
    factor_rows = []
    for key in ["savings_per_mw", "payback_period", "impact_on_opex",
                "water_savings_pct", "load_growth"]:
        factor_rows.append(factor_display[key])

    add_table(slide, right_x, Inches(1.5), right_w, Inches(2.2),
              factor_headers, factor_rows, font_size=9)

    # Normalization formula
    add_textbox(slide, right_x, Inches(3.9), right_w, Inches(0.4),
                "Normalization Formula", font_size=14, bold=True, color=DARK_BLUE)

    tf = add_textbox(slide, right_x, Inches(4.3), right_w, Inches(1.5),
                     "Standard:  score = clip( (value \u2013 min) / (max \u2013 min) \u00d7 100,  0, 100 )",
                     font_size=11, color=DARK_GRAY, font_name="Consolas")
    add_para(tf, "Inverted:   score = clip( (1 \u2013 (value \u2013 min) / (max \u2013 min)) \u00d7 100,  0, 100 )",
             font_size=11, color=DARK_GRAY, font_name="Consolas", space_before=Pt(4))
    add_para(tf, "", font_size=6)
    add_para(tf, "Composite = \u03a3 (factor_score \u00d7 weight)  \u2192  0\u2013100 final score",
             font_size=12, bold=True, color=MEDIUM_BLUE, font_name="Consolas", space_before=Pt(8))

    # Input summary mini-table
    add_textbox(slide, right_x, Inches(5.4), right_w, Inches(0.4),
                "Input Summary (9 parameters)", font_size=14, bold=True, color=DARK_BLUE)

    inp_headers = ["Input", "Default", "Source", "Status"]
    inp_rows = []
    for inp in inputs_spec["inputs"]:
        default_val = inp["default"]
        if default_val is None:
            default_str = "from CSV"
        elif isinstance(default_val, float) and default_val < 1:
            default_str = f"{default_val:.0%}"
        elif isinstance(default_val, (int, float)) and default_val >= 1000:
            default_str = f"${default_val:,.0f}"
        else:
            default_str = str(default_val)
        status = inp["real_or_estimated"]
        if status == "REAL":
            status_str = "\u2705 Real"
        elif status == "REAL_STATE_AVG":
            status_str = "\u2705 Real (state avg)"
        elif status == "PARAM":
            status_str = "\u2699 Vigilent param"
        else:
            status_str = "\u26a0 Estimated"
        inp_rows.append((inp["name"], default_str, inp["source"], status_str))

    add_table(slide, right_x, Inches(5.8), right_w, Inches(1.5),
              inp_headers, inp_rows, font_size=8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: WORKED EXAMPLE + RANKINGS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_2(prs, all_dcs, example_dc):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, WHITE)

    # Title bar
    add_textbox(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.8),
                f"  Worked Example: {example_dc['name']} ({example_dc['city']}, "
                f"{example_dc['state']}) \u2014 {example_dc['size_mw']:.0f} MW",
                font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    title_shape = slide.shapes[0]
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE

    # --- Recompute step-by-step with real values ---
    mw = example_dc["size_mw"]
    pue = example_dc["baseline_pue"]
    price = example_dc["electricity_price"]
    growth = example_dc["load_growth_rate"]
    opex_pct = example_dc["energy_pct_opex"]
    invest = 1_500_000
    e_red = 0.10
    w_red = 0.05

    total_mw = mw * pue
    annual_kwh = total_mw * 1000 * 8760 * (1 + growth)
    annual_cost = annual_kwh * price
    savings = annual_cost * e_red
    sav_per_mw = savings / mw
    payback = invest / savings
    opex_impact = e_red * opex_pct

    # Factor scores
    result = compute_score(
        dc_size_mw=mw, baseline_pue=pue, electricity_price=price,
        load_growth_rate=growth, energy_pct_opex=opex_pct,
        investment_cost=invest, energy_reduction_pct=e_red,
        water_reduction_pct=w_red,
    )
    fs = result["factor_scores"]
    composite = result["composite_score"]

    # --- Left column: Step-by-step calculation ---
    left_x = Inches(0.4)
    col_w = Inches(6.2)

    tf = add_textbox(slide, left_x, Inches(0.95), col_w, Inches(0.3),
                     "Step-by-Step Calculation", font_size=14, bold=True, color=DARK_BLUE)

    steps = [
        ("1. Total MW (with cooling overhead)",
         f"{mw:.1f} MW \u00d7 {pue} PUE = {total_mw:.1f} MW"),
        ("2. Annual Energy Consumption",
         f"{total_mw:.1f} MW \u00d7 1,000 \u00d7 8,760 hrs \u00d7 (1 + {growth:.0%}) = {annual_kwh:,.0f} kWh"),
        ("3. Annual Energy Cost",
         f"{annual_kwh:,.0f} kWh \u00d7 ${price}/kWh = ${annual_cost:,.2f}"),
        ("4. Estimated Savings (10% reduction)",
         f"${annual_cost:,.2f} \u00d7 10% = ${savings:,.2f}"),
        ("5. Savings per MW",
         f"${savings:,.2f} / {mw:.1f} MW = ${sav_per_mw:,.2f}/MW"),
        ("6. Payback Period",
         f"$1,500,000 / ${savings:,.2f} = {payback:.3f} years"),
        ("7. OPEX Impact",
         f"10% \u00d7 40% = {opex_impact:.0%}"),
    ]

    y = Inches(1.3)
    for label, formula in steps:
        tf = add_textbox(slide, left_x, y, col_w, Inches(0.40),
                         label, font_size=10, bold=True, color=MEDIUM_BLUE)
        add_para(tf, f"  {formula}", font_size=10, color=DARK_GRAY,
                 font_name="Consolas", space_before=Pt(1))
        y += Inches(0.42)

    # Factor scoring summary
    y += Inches(0.1)
    tf = add_textbox(slide, left_x, y, col_w, Inches(0.3),
                     "Factor Scores \u2192 Composite", font_size=14, bold=True, color=DARK_BLUE)
    y += Inches(0.35)

    factor_labels = [
        ("Savings/MW", "savings_per_mw", 0.35),
        ("Payback (inv)", "payback_period", 0.25),
        ("OPEX Impact", "impact_on_opex", 0.20),
        ("Water Savings", "water_savings_pct", 0.10),
        ("Load Growth", "load_growth", 0.10),
    ]

    weighted_parts = []
    for label, key, weight in factor_labels:
        score = fs[key]
        weighted = score * weight
        weighted_parts.append(f"{weighted:.1f}")

    tf = add_textbox(slide, left_x, y, col_w, Inches(0.3), "", font_size=1)
    for label, key, weight in factor_labels:
        score = fs[key]
        weighted = score * weight
        add_para(tf, f"{label}: {score:.1f}/100 \u00d7 {weight:.0%} = {weighted:.1f}",
                 font_size=10, color=DARK_GRAY, font_name="Consolas", space_before=Pt(1))

    add_para(tf, "", font_size=4)
    add_para(tf, f"Composite = {' + '.join(weighted_parts)} = {composite:.2f}/100  \u2192  {example_dc['classification']}",
             font_size=12, bold=True, color=GREEN, font_name="Consolas", space_before=Pt(4))

    # --- Right column: Full rankings table ---
    right_x = Inches(7.0)
    right_w = Inches(5.9)

    add_textbox(slide, right_x, Inches(0.95), right_w, Inches(0.3),
                f"All {len(all_dcs)} Data Centers \u2014 Ranked by Score",
                font_size=14, bold=True, color=DARK_BLUE)

    rank_headers = ["#", "Name", "State", "MW", "Score", "Class", "Payback"]
    rank_rows = []
    for i, dc in enumerate(all_dcs):
        rank_rows.append((
            str(i + 1),
            dc["name"],
            dc["state"],
            f"{dc['size_mw']:.1f}",
            f"{dc['composite_score']:.1f}",
            dc["classification"],
            f"{dc['payback_years']:.2f} yr",
        ))

    # Dynamic row height based on DC count
    table_height = min(Inches(6.0), Inches(0.22 * (len(rank_rows) + 1)))
    add_table(slide, right_x, Inches(1.3), right_w, table_height,
              rank_headers, rank_rows, font_size=8)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    print("Generating client walkthrough deck...")

    # Load data
    with open(SCORED_JSON) as f:
        all_dcs = json.load(f)
    with open(INPUTS_JSON) as f:
        inputs_spec = json.load(f)

    # Find example DC
    example_dc = next((dc for dc in all_dcs if dc["name"] == EXAMPLE_DC_NAME), all_dcs[0])

    # Build deck
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs, all_dcs, inputs_spec)
    build_slide_2(prs, all_dcs, example_dc)

    # Save
    os.makedirs("output", exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Example DC: {example_dc['name']} ({example_dc['city']}, {example_dc['state']})")
    print(f"  Composite Score: {example_dc['composite_score']}")
    print(f"  Total DCs ranked: {len(all_dcs)}")


if __name__ == "__main__":
    main()
